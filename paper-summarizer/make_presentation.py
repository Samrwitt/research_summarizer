
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def add_title_slide(prs, title_text, subtitle_text, author_text):
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = title_text
    subtitle.text = f"{subtitle_text}\n\n{author_text}"
    
def add_content_slide(prs, title_text, content_points):
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = title_text
    tf = body.text_frame
    
    for i, point in enumerate(content_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(20)
        p.space_after = Pt(10)

def create_presentation():
    prs = Presentation()
    
    # 1. Title
    add_title_slide(
        prs, 
        "Research Summarizer AI", 
        "An Advanced NLP System for Scientific Literature", 
        "Project Overview"
    )
    
    # 2. Problem Statement
    add_content_slide(
        prs,
        "The Problem: Information Overload",
        [
            "Exponential growth in scientific publications makes it impossible to keep up.",
            "Technical density of papers creates a high barrier to entry for students.",
            "Existing tools are either too generic (ChatGPT) or too simple (keyword matching).",
            "Need for a system that understands structure (Methods, Results) not just words."
        ]
    )
    
    # 3. Project Objectives
    add_content_slide(
        prs,
        "Project Objectives",
        [
            "Develop an automatic summarization system tailored for research papers.",
            "Handle long documents (10k+ tokens) via Section-Aware Preprocessing and Longformer models.",
            "Implement 'Serious NLP': ROUGE metric evaluation and hybrid extraction strategies.",
            "Create a premium, research-grade user experience (Markdown/DOCX exports)."
        ]
    )
    
    # 4. Methodology: The Pipeline
    add_content_slide(
        prs,
        "Methodology: End-to-End Pipeline",
        [
            "1. Ingestion: Automatic fetching from arXiv IDs or PDF uploads.",
            "2. Preprocessing: Regex-based cleaning and Section Identification (Intro vs Methods).",
            "3. Extractive Summarization: TF-IDF ranking to identify key sentences as a baseline.",
            "4. Abstractive Summarization: Using 'Longformer Encoder-Decoder' (LED), a sparse-attention Transformer designed for long texts.",
            "5. Evaluation: Automated ROUGE scoring against reference abstracts."
        ]
    )
    
    # 5. Technical Architecture
    add_content_slide(
        prs,
        "Technical Stack",
        [
            "Language: Python 3.10+",
            "Model: HuggingFace Transformers (AllenAI LED-base-16384)",
            "NLP Libs: NLTK (Tokenization), Scikit-Learn (TF-IDF), ROUGE-Score",
            "Interface: Streamlit (Custom CSS Premium UI)",
            "Infrastructure: Local-First (Privacy focused, offline capable)"
        ]
    )
    
    # 6. Why This Matters (vs GPT)
    add_content_slide(
        prs,
        "Competitive Advantage: Why not just GPT?",
        [
            "Specialization: Trained specifically on scientific discourse, less 'chatty' fluff.",
            "Privacy: Runs locally. Critical for confidential R&D or legal documents.",
            "Workflow Integration: Batch processing, DOCX export, and structured metadata extraction.",
            "Verifiability: Built-in quantitative evaluation (ROUGE scores) to prove accuracy."
        ]
    )
    
    # 7. Future Work
    add_content_slide(
        prs,
        "Future Enhancements",
        [
            "Batch Processing: Summarize 100+ PDFs at once into a CSV review matrix.",
            "Q&A Module: RAG (Retrieval Augmented Generation) for chatting with the paper.",
            "Citation Graphing: Visualizing connections between references.",
            "Fine-tuning: Training the LED model on specific domain data (e.g., BioMedical)."
        ]
    )
    
    count = len(prs.slides)
    output_path = "Research_Summarizer_Project.pptx"
    prs.save(output_path)
    print(f"Presentation saved to {output_path} with {count} slides.")

if __name__ == "__main__":
    create_presentation()
