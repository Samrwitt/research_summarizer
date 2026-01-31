
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def add_title_slide(prs, title_text, subtitle_text, author_text):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = title_text
    subtitle.text = f"{subtitle_text}\n\n{author_text}"
    
def add_content_slide(prs, title_text, content_points):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = title_text
    tf = body.text_frame
    tf.word_wrap = True
    
    for i, point in enumerate(content_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        # Handle indentation for sub-bullets
        if point.strip().startswith("•"):
            p.text = point.replace("•", "").strip()
            p.level = 1
            p.font.size = Pt(18)
        else:
            p.text = point
            p.level = 0
            p.font.size = Pt(20)
        
        p.space_after = Pt(8)

def create_presentation():
    prs = Presentation()
    
    # ===== SLIDE 1: Title =====
    add_title_slide(
        prs, 
        "Research Summarizer AI", 
        "An NLP System for Intelligent Scientific Document Summarization", 
        "Technical Deep Dive: Models, Methods & Process"
    )
    
    # ===== SLIDE 2: System Overview =====
    add_content_slide(
        prs,
        "System Architecture Overview",
        [
            "End-to-End AI Pipeline for Research Paper Processing",
            "• Multi-source ingestion (arXiv, PDF, Text)",
            "• Section-aware intelligent preprocessing",
            "• Hybrid extractive-abstractive summarization",
            "• Semantic analysis with KeyBERT + Ollama LLM",
            "• Interactive RAG-based Q&A system",
            "• ROUGE-based quantitative evaluation"
        ]
    )
    
    # ===== SLIDE 3: Complete Process Flow =====
    add_content_slide(
        prs,
        "Complete Summarization Process (Step-by-Step)",
        [
            "1. INGESTION → Fetch from arXiv API or extract from PDF",
            "2. CLEANING → Remove hyphenation, references, normalize text",
            "3. SECTION EXTRACTION → Identify Intro/Methods/Results/Conclusion",
            "4. FOCUS BUILDING → Prioritize Abstract → Conclusion → Results",
            "5. CHUNKING → Split into 3000-char overlapping segments",
            "6. SUMMARIZATION → Extract key sentences + Abstractive generation",
            "7. POST-PROCESSING → Generate bullets, export to MD/DOCX"
        ]
    )
    
    # ===== SLIDE 4: Ingestion Methods =====
    add_content_slide(
        prs,
        "Step 1: Multi-Source Ingestion",
        [
            "Three ingestion pathways implemented:",
            "• arXiv API: Direct fetch via paper ID with metadata",
            "• PDF Upload: pdfplumber → pypdf → pdftotext fallback chain",
            "• Text Files: Direct UTF-8 ingestion",
            "",
            "Why multiple methods?",
            "• Maximizes compatibility across file formats",
            "• Ensures robustness when one parser fails"
        ]
    )
    
    # ===== SLIDE 5: Preprocessing Deep Dive =====
    add_content_slide(
        prs,
        "Step 2-4: Intelligent Preprocessing",
        [
            "Text Cleaning (Regex-based):",
            "• Remove PDF artifacts (hyphenation: 'process-ing' → 'processing')",
            "• Strip references section (last 20% heuristic search)",
            "• Normalize whitespace and flatten structure",
            "",
            "Section-Aware Extraction:",
            "• Regex patterns identify Abstract/Intro/Methods/Results/Conclusion",
            "• Build focus text prioritizing high-signal sections",
            "• Methods limited to 2000 chars (avoid equation overload)"
        ]
    )
    
    # ===== SLIDE 6: Chunking Strategy =====
    add_content_slide(
        prs,
        "Step 5: Chunking for Long Context",
        [
            "Challenge: Papers often exceed 10,000 tokens",
            "",
            "Solution: Sliding window chunking",
            "• Chunk size: 3000 characters (~750 tokens)",
            "• Overlap: 200 characters (preserve context continuity)",
            "• Smart breaking at word boundaries",
            "",
            "Why not process full text?",
            "• Even LED has practical memory limits at inference",
            "• Chunking enables parallel processing"
        ]
    )
    
    # ===== SLIDE 7: Extractive Model =====
    add_content_slide(
        prs,
        "Model 1: Extractive Summarization (TF-IDF)",
        [
            "Method: Sentence ranking via sklearn TfidfVectorizer",
            "• Tokenize text into sentences (NLTK)",
            "• Compute TF-IDF scores for each sentence",
            "• Rank by sum of term weights, select top-N",
            "",
            "Why TF-IDF vs BERT embeddings?",
            "• 100x faster (no GPU required)",
            "• Interpretable scores (debugging/explainability)",
            "• Sufficient for identifying salient sentences"
        ]
    )
    
    # ===== SLIDE 8: Abstractive Model =====
    add_content_slide(
        prs,
        "Model 2: Abstractive Summarization (LED)",
        [
            "Model: AllenAI Longformer Encoder-Decoder (LED-base-16384)",
            "• Transformer with sparse attention (global + local)",
            "• Pre-trained on arXiv + PubMed scientific papers",
            "• Context window: 16,384 tokens (~4x longer than BERT)",
            "",
            "Why LED specifically?",
            "• Designed for long documents (attention is O(n) not O(n²))",
            "• Domain-specific pre-training on scientific text",
            "• Generates fluent, coherent summaries (not just extraction)"
        ]
    )
    
    # ===== SLIDE 9: Alternative DistilBART =====
    add_content_slide(
        prs,
        "Fallback: DistilBART-CNN (Why We Offer This)",
        [
            "Model: sshleifer/distilbart-cnn-12-6",
            "• Distilled BART model (smaller, faster)",
            "• Trained on CNN/DailyMail news dataset",
            "",
            "When to use DistilBART instead of LED?",
            "• Papers \u003c 1024 tokens (short abstracts/letters)",
            "• Limited GPU memory (LED requires ~6GB, DistilBART ~2GB)",
            "• Faster inference needed (3x speedup)",
            "",
            "Trade-off: Less context, more generic (news vs science)"
        ]
    )
    
    # ===== SLIDE 10: Hybrid Mode =====
    add_content_slide(
        prs,
        "Hybrid Mode: Best of Both Worlds",
        [
            "Problem: LED is slow for very long papers (10k+ tokens)",
            "",
            "Solution: Two-stage hybrid pipeline",
            "1. Extractive pre-filtering (TF-IDF) → Select top 50% sentences",
            "2. Abstractive refinement (LED) → Generate from filtered text",
            "",
            "Benefits:",
            "• 50% faster processing time",
            "• Maintains global context (extractive sees full text)",
            "• Produces fluent output (abstractive final pass)"
        ]
    )
    
    # ===== SLIDE 11: Why Not Other Models? =====
    add_content_slide(
        prs,
        "Model Selection: Why Not Alternatives?",
        [
            "Why not GPT-4 / ChatGPT?",
            "• Privacy: No data leaves local machine",
            "• Cost: Zero API fees",
            "• Reproducibility: Deterministic outputs",
            "",
            "Why not T5 or PEGASUS?",
            "• T5: Generic pre-training (not scientific domain)",
            "• PEGASUS: 512 token limit (same as BERT)",
            "• LED: Explicitly designed for long scientific docs"
        ]
    )
    
    # ===== SLIDE 12: Semantic Analysis =====
    add_content_slide(
        prs,
        "Semantic Insights: KeyBERT + Ollama",
        [
            "Component 1: KeyBERT (Keyword Extraction)",
            "• Uses sentence-transformers (all-MiniLM-L6-v2)",
            "• Extracts top 10 statistically significant keywords",
            "",
            "Component 2: Ollama LLM (llama3.2)",
            "• Generates natural language analysis from keywords",
            "• Runs locally via Ollama REST API",
            "• Fallback to 'keywords only' if Ollama offline",
            "",
            "Why this combination? Fast extraction + contextual understanding"
        ]
    )
    
    # ===== SLIDE 13: RAG System =====
    add_content_slide(
        prs,
        "Interactive Q&A: RAG Implementation",
        [
            "Retrieval-Augmented Generation for paper chat",
            "",
            "Pipeline:",
            "1. Embed all sentences (SentenceTransformer: all-MiniLM-L6-v2)",
            "2. Embed user question",
            "3. Cosine similarity search → retrieve top-5 sentences",
            "4. Pass context + question to Ollama (llama3.2)",
            "5. Fallback to DistilBERT-SQuAD if Ollama unavailable",
            "",
            "Why this architecture? Local-first, privacy-preserving Q&A"
        ]
    )
    
    # ===== SLIDE 14: Evaluation Methodology =====
    add_content_slide(
        prs,
        "Quantitative Evaluation: ROUGE Metrics",
        [
            "ROUGE (Recall-Oriented Understudy for Gisting Evaluation)",
            "• ROUGE-1: Unigram overlap (keyword coverage)",
            "• ROUGE-2: Bigram overlap (phrase preservation)",
            "• ROUGE-L: Longest common subsequence (structure similarity)",
            "",
            "Benchmark: Compare generated summary vs. author's abstract",
            "",
            "Why ROUGE? Industry-standard for summarization evaluation",
            "Used in research papers (e.g., BART, PEGASUS, LED validation)"
        ]
    )
    
    # ===== SLIDE 15: Technical Stack Summary =====
    add_content_slide(
        prs,
        "Complete Technical Stack",
        [
            "Ingestion: pdfplumber, pypdf, BeautifulSoup, requests",
            "NLP Processing: NLTK (tokenization), regex (cleaning)",
            "Extractive: scikit-learn (TF-IDF)",
            "Abstractive: HuggingFace Transformers (LED, DistilBART)",
            "Embeddings: sentence-transformers (MiniLM)",
            "LLM: Ollama (llama3.2 via REST API)",
            "Evaluation: rouge-score",
            "UI: Streamlit with custom CSS"
        ]
    )
    
    # ===== SLIDE 16: Why This Matters =====
    add_content_slide(
        prs,
        "Real-World Impact & Use Cases",
        [
            "Who benefits from this system?",
            "• Graduate students reviewing 100+ papers for literature reviews",
            "• R&D teams in pharma/biotech (confidential documents)",
            "• Legal professionals analyzing patent documents",
            "",
            "Key advantages:",
            "• 10x faster than manual reading",
            "• Privacy-preserving (no cloud APIs)",
            "• Quantitatively validated (ROUGE scores)"
        ]
    )
    
    # ===== SLIDE 17: Future Work =====
    add_content_slide(
        prs,
        "Future Enhancements",
        [
            "Model improvements:",
            "• Fine-tune LED on domain-specific corpora (e.g., biomedical)",
            "• Experiment with newer models (LongT5, LED-large)",
            "",
            "Feature additions:",
            "• Batch processing (100+ PDFs → CSV matrix)",
            "• Citation graph visualization",
            "• Multi-document summarization for literature reviews",
            "• Figure/table extraction and captioning"
        ]
    )
    
    output_path = "Research_Summarizer_Project.pptx"
    prs.save(output_path)
    print(f"✓ Presentation saved to {output_path} with {len(prs.slides)} slides.")

if __name__ == "__main__":
    create_presentation()
